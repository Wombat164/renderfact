"""
Tests for roundtrip/provenance.py -- D11 part 2 (chunk 4.1): hidden provenance
metadata embedded in every rendered editable-document artifact -- DOCX, XLSX,
PPTX (not SVG/PNG -- visual/diagram artifacts -- and not PDF -- a flattened,
non-round-trippable archival format).

Covers: embed/extract round-trips against REAL files in all three formats
(created via python-docx/openpyxl/python-pptx, not a mock) -- proving the
SAME mechanism (OOXML dc:identifier) genuinely works across all three, not
just asserted to; an unsupported extension raises a clear error; a file with
no renderfact provenance extracts as None; re-embedding overwrites rather
than accumulates; the embedded field is genuinely invisible in the rendered
body; tool_version() resolves against this repo's real git state; and the
render.py CLI dispatch (embed then extract) works end-to-end including the
source-UID side effect on the .md file.
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import pytest
from docx import Document

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT))

from roundtrip import provenance  # noqa: E402
from roundtrip import source_uid  # noqa: E402

RENDER_PY = REPO_ROOT / "render.py"


def _make_docx(tmp_path: Path, body_text: str = "Hello, world.") -> Path:
    doc = Document()
    doc.add_paragraph(body_text)
    path = tmp_path / "rendered.docx"
    doc.save(str(path))
    return path


def _make_xlsx(tmp_path: Path, cell_text: str = "Hello, world.") -> Path:
    import openpyxl

    wb = openpyxl.Workbook()
    wb.active["A1"] = cell_text
    path = tmp_path / "rendered.xlsx"
    wb.save(str(path))
    return path


def _make_pptx(tmp_path: Path, body_text: str = "Hello, world.") -> Path:
    import pptx

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text_frame.text = body_text
    path = tmp_path / "rendered.pptx"
    prs.save(str(path))
    return path


_MAKERS = {".docx": _make_docx, ".xlsx": _make_xlsx, ".pptx": _make_pptx}


@pytest.mark.parametrize("suffix", [".docx", ".xlsx", ".pptx"])
def test_embed_then_extract_round_trips_across_all_three_formats(tmp_path, suffix):
    artifact_path = _MAKERS[suffix](tmp_path)
    prov = provenance.Provenance(
        source_uid="11111111-1111-1111-1111-111111111111",
        source_version="abc123def4567890",
        rendered_at="2026-07-03T12:00:00Z",
        tool_version="78a3110",
    )
    provenance.embed(artifact_path, prov)

    extracted = provenance.extract(artifact_path)
    assert extracted == prov


def test_embed_rejects_unsupported_extension(tmp_path):
    svg_path = tmp_path / "diagram.svg"
    svg_path.write_text("<svg></svg>", encoding="utf-8")
    prov = provenance.Provenance("uid", "v", "2026-07-03T00:00:00Z", "tool")

    with pytest.raises(ValueError, match="unsupported artifact type '.svg'"):
        provenance.embed(svg_path, prov)


def test_embed_then_extract_round_trips(tmp_path):
    docx_path = _make_docx(tmp_path)
    prov = provenance.Provenance(
        source_uid="11111111-1111-1111-1111-111111111111",
        source_version="abc123def4567890",
        rendered_at="2026-07-03T12:00:00Z",
        tool_version="78a3110",
    )
    provenance.embed(docx_path, prov)

    extracted = provenance.extract(docx_path)
    assert extracted == prov


def test_extract_returns_none_for_docx_with_no_provenance(tmp_path):
    docx_path = _make_docx(tmp_path)
    assert provenance.extract(docx_path) is None


def test_extract_returns_none_for_unrelated_identifier_value(tmp_path):
    # A DOCX whose dc:identifier was set by something else entirely (not
    # renderfact) must not be misread as provenance.
    docx_path = _make_docx(tmp_path)
    doc = Document(str(docx_path))
    doc.core_properties.identifier = "ISBN-978-0-13-468599-1"
    doc.save(str(docx_path))

    assert provenance.extract(docx_path) is None


def test_reembedding_overwrites_not_accumulates(tmp_path):
    docx_path = _make_docx(tmp_path)
    first = provenance.Provenance("uid-1", "v1", "2026-07-01T00:00:00Z", "tool-1")
    second = provenance.Provenance("uid-2", "v2", "2026-07-03T00:00:00Z", "tool-2")

    provenance.embed(docx_path, first)
    provenance.embed(docx_path, second)

    extracted = provenance.extract(docx_path)
    assert extracted == second


def test_provenance_is_not_visible_in_rendered_body(tmp_path):
    docx_path = _make_docx(tmp_path, body_text="Only this text should be visible.")
    prov = provenance.Provenance("uid", "v", "2026-07-03T00:00:00Z", "tool")
    provenance.embed(docx_path, prov)

    doc = Document(str(docx_path))
    body_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Only this text should be visible." in body_text
    assert "uid" not in body_text
    assert "renderfact" not in body_text


def test_build_provenance_raises_clean_error_for_missing_source(tmp_path):
    # Previously a bare FileNotFoundError traceback -- must now be a clear,
    # actionable ProvenanceError pointing at adopt() instead.
    missing = tmp_path / "does-not-exist.md"
    with pytest.raises(provenance.ProvenanceError, match="use 'adopt' instead"):
        provenance.build_provenance(missing)


def test_adopt_bootstraps_a_stub_source_and_embeds_provenance(tmp_path):
    docx_path = _make_docx(tmp_path, body_text="Hand-authored content, never rendered by renderfact.")
    new_source = tmp_path / "adopted-source.md"
    assert not new_source.exists()

    prov = provenance.adopt(docx_path, new_source)

    assert new_source.exists()
    stub = new_source.read_text(encoding="utf-8")
    assert "origin: adopted-external-draft" in stub
    assert "adopted_at:" in stub

    extracted = provenance.extract(docx_path)
    assert extracted == prov
    assert prov.source_uid in stub


def test_adopt_refuses_when_artifact_already_has_provenance(tmp_path):
    docx_path = _make_docx(tmp_path)
    existing_source = tmp_path / "existing.md"
    existing_source.write_text("---\ntitle: X\n---\n\nBody.\n", encoding="utf-8")
    provenance.embed(docx_path, provenance.build_provenance(existing_source))

    new_source = tmp_path / "should-not-be-created.md"
    with pytest.raises(provenance.ProvenanceError, match="already carries renderfact provenance"):
        provenance.adopt(docx_path, new_source)
    assert not new_source.exists()


def test_adopt_refuses_when_source_path_already_exists(tmp_path):
    docx_path = _make_docx(tmp_path)  # no provenance yet
    already_there = tmp_path / "already-there.md"
    already_there.write_text("pre-existing content\n", encoding="utf-8")

    with pytest.raises(provenance.ProvenanceError, match="already exists"):
        provenance.adopt(docx_path, already_there)
    # the pre-existing source content must survive untouched
    assert already_there.read_text(encoding="utf-8") == "pre-existing content\n"
    assert provenance.extract(docx_path) is None


def test_retarget_carries_identity_and_content_version_onto_new_format(tmp_path):
    docx_path = _make_docx(tmp_path)
    source = tmp_path / "source.md"
    source.write_text("---\ntitle: X\n---\n\nBody.\n", encoding="utf-8")
    old_prov = provenance.build_provenance(source)
    provenance.embed(docx_path, old_prov)

    xlsx_path = _make_xlsx(tmp_path)  # "already produced" new-format file, no provenance yet

    new_prov = provenance.retarget(docx_path, xlsx_path)

    assert new_prov.source_uid == old_prov.source_uid
    assert new_prov.source_version == old_prov.source_version
    # a genuinely new physical artifact -- its own rendered_at, not a copy of the old one
    assert provenance.extract(xlsx_path) == new_prov
    # the old artifact is untouched
    assert provenance.extract(docx_path) == old_prov


def test_retarget_refuses_when_old_artifact_has_no_provenance(tmp_path):
    docx_path = _make_docx(tmp_path)  # never embedded
    xlsx_path = _make_xlsx(tmp_path)

    with pytest.raises(provenance.ProvenanceError, match="use 'adopt' on the new artifact"):
        provenance.retarget(docx_path, xlsx_path)


def test_retarget_refuses_when_new_artifact_does_not_exist_yet(tmp_path):
    docx_path = _make_docx(tmp_path)
    source = tmp_path / "source.md"
    source.write_text("---\ntitle: X\n---\n\nBody.\n", encoding="utf-8")
    provenance.embed(docx_path, provenance.build_provenance(source))

    not_yet_rendered = tmp_path / "not-yet-rendered.xlsx"
    with pytest.raises(provenance.ProvenanceError, match="does not exist yet"):
        provenance.retarget(docx_path, not_yet_rendered)


def test_retarget_refuses_to_overwrite_a_differently_tracked_new_artifact(tmp_path):
    docx_path = _make_docx(tmp_path)
    source_a = tmp_path / "source-a.md"
    source_a.write_text("---\ntitle: A\n---\n\nBody A.\n", encoding="utf-8")
    provenance.embed(docx_path, provenance.build_provenance(source_a))

    xlsx_path = _make_xlsx(tmp_path)
    source_b = tmp_path / "source-b.md"
    source_b.write_text("---\ntitle: B\n---\n\nBody B.\n", encoding="utf-8")
    unrelated_prov = provenance.build_provenance(source_b)
    provenance.embed(xlsx_path, unrelated_prov)  # xlsx already tracks a DIFFERENT document

    with pytest.raises(provenance.ProvenanceError, match="already carries DIFFERENT renderfact provenance"):
        provenance.retarget(docx_path, xlsx_path)
    # the unrelated artifact's own provenance must survive untouched
    assert provenance.extract(xlsx_path) == unrelated_prov


def test_retarget_is_a_noop_refusal_when_new_artifact_already_matches(tmp_path):
    # Retargeting onto a file that ALREADY carries the SAME source_uid (e.g. a
    # re-run) is allowed -- only a DIFFERENT identity is refused.
    docx_path = _make_docx(tmp_path)
    source = tmp_path / "source.md"
    source.write_text("---\ntitle: X\n---\n\nBody.\n", encoding="utf-8")
    prov = provenance.build_provenance(source)
    provenance.embed(docx_path, prov)

    xlsx_path = _make_xlsx(tmp_path)
    provenance.embed(xlsx_path, prov)  # already tagged with the SAME source_uid

    result = provenance.retarget(docx_path, xlsx_path)
    assert result.source_uid == prov.source_uid


def test_tool_version_resolves_against_real_repo_git_state():
    import re

    v = provenance.tool_version()
    assert v != "unknown"
    # Either a semver tag (the primary source once releases exist, per OQ7:
    # e.g. v0.1.0, possibly with -N-gHASH suffix) or a bare short commit hash
    # (dev-build fallback). The original >=7 length assertion encoded the
    # pre-tag world and failed the moment a v0.1.0 tag existed.
    assert re.match(r"^v\d+\.\d+", v) or re.match(r"^[0-9a-f]{7,}", v), v


def test_now_iso_is_utc_z_suffixed_and_parseable():
    from datetime import datetime

    stamp = provenance.now_iso()
    assert stamp.endswith("Z")
    datetime.strptime(stamp, "%Y-%m-%dT%H:%M:%SZ")  # raises if malformed


def test_cli_embed_then_extract_end_to_end(tmp_path):
    md_path = tmp_path / "source.md"
    md_path.write_text("---\ntitle: Test\n---\n\nSome content.\n", encoding="utf-8")
    docx_path = _make_docx(tmp_path)

    embed_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "embed", str(docx_path), "--source", str(md_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert embed_result.returncode == 0, embed_result.stderr
    assert "embedded provenance" in embed_result.stdout

    # the source-UID side effect actually happened on the .md file
    md_content = md_path.read_text(encoding="utf-8")
    assert "renderfact_uid:" in md_content

    extract_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "extract", str(docx_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert extract_result.returncode == 0, extract_result.stderr
    payload = json.loads(extract_result.stdout)
    assert payload["source_uid"] in md_content
    assert len(payload["source_version"]) == 16


def test_cli_extract_reports_absence_with_nonzero_exit(tmp_path):
    docx_path = _make_docx(tmp_path)
    result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "extract", str(docx_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert result.returncode == 1
    assert "no renderfact provenance found" in result.stderr


def test_cli_embed_missing_source_gives_clean_error_not_a_traceback(tmp_path):
    docx_path = _make_docx(tmp_path)
    missing_source = tmp_path / "does-not-exist.md"
    result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "embed", str(docx_path), "--source", str(missing_source)],
        capture_output=True, text=True, timeout=30,
    )
    assert result.returncode == 1
    assert "Traceback" not in result.stderr
    assert "use 'adopt' instead" in result.stderr


def test_cli_adopt_then_extract_end_to_end(tmp_path):
    docx_path = _make_docx(tmp_path, body_text="Drafted directly in Word, no renderfact history.")
    new_source = tmp_path / "adopted.md"

    adopt_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "adopt", str(docx_path), "--source", str(new_source)],
        capture_output=True, text=True, timeout=30,
    )
    assert adopt_result.returncode == 0, adopt_result.stderr
    assert "adopted" in adopt_result.stdout
    assert new_source.exists()
    assert "origin: adopted-external-draft" in new_source.read_text(encoding="utf-8")

    extract_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "extract", str(docx_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert extract_result.returncode == 0, extract_result.stderr
    payload = json.loads(extract_result.stdout)
    assert payload["source_uid"] in new_source.read_text(encoding="utf-8")


def test_cli_retarget_end_to_end(tmp_path):
    docx_path = _make_docx(tmp_path)
    source = tmp_path / "source.md"
    source.write_text("---\ntitle: X\n---\n\nBody.\n", encoding="utf-8")
    embed_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "embed", str(docx_path), "--source", str(source)],
        capture_output=True, text=True, timeout=30,
    )
    assert embed_result.returncode == 0, embed_result.stderr

    xlsx_path = _make_xlsx(tmp_path)  # already-produced new-format file, no provenance yet
    retarget_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "retarget", str(docx_path), str(xlsx_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert retarget_result.returncode == 0, retarget_result.stderr
    assert "retargeted" in retarget_result.stdout

    extract_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "extract", str(xlsx_path)],
        capture_output=True, text=True, timeout=30,
    )
    payload = json.loads(extract_result.stdout)
    assert payload["source_uid"] in source.read_text(encoding="utf-8")


def test_cli_embed_then_extract_end_to_end_for_xlsx(tmp_path):
    # Same CLI dispatch path as the DOCX test above, proving the multi-format
    # support is real at the render.py entry point too, not just in the
    # directly-called provenance.embed()/extract() functions.
    md_path = tmp_path / "source.md"
    md_path.write_text("---\ntitle: Test\n---\n\nSome content.\n", encoding="utf-8")
    xlsx_path = _make_xlsx(tmp_path)

    embed_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "embed", str(xlsx_path), "--source", str(md_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert embed_result.returncode == 0, embed_result.stderr

    extract_result = subprocess.run(
        [sys.executable, str(RENDER_PY), "provenance", "extract", str(xlsx_path)],
        capture_output=True, text=True, timeout=30,
    )
    assert extract_result.returncode == 0, extract_result.stderr
    payload = json.loads(extract_result.stdout)
    assert payload["source_uid"] in md_path.read_text(encoding="utf-8")
